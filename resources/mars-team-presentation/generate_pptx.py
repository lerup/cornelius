#!/usr/bin/env python3
"""Generate Mars Team Talk PPTX from slide content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# Colors
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x80, 0x80, 0x80)
LIGHT_GREY = RGBColor(0xAA, 0xAA, 0xAA)
VERDANT = RGBColor(0xB3, 0xBC, 0xB5)
VERDANT_DARK = RGBColor(0x5A, 0x6B, 0x5D)
AZURE = RGBColor(0xB4, 0xBA, 0xCC)
AZURE_DARK = RGBColor(0x8A, 0x90, 0xA8)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
BG_GREY = RGBColor(0xF0, 0xF0, 0xF0)
CARD_VERDANT = RGBColor(0xE8, 0xED, 0xE9)
CARD_AZURE = RGBColor(0xE8, 0xEA, 0xF0)

# Dimensions
SLIDE_W = Inches(13.333)  # 1920px at 144dpi = 13.333in (widescreen)
SLIDE_H = Inches(7.5)     # 1080px at 144dpi = 7.5in
MARGIN = Inches(0.9)
MARGIN_SM = Inches(0.6)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=BLACK, alignment=PP_ALIGN.LEFT, font_name='Calibri',
                line_spacing=None):
    """Add a text box with single-style text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox


def add_rich_textbox(slide, left, top, width, height, runs, alignment=PP_ALIGN.LEFT,
                     line_spacing=None):
    """Add a text box with mixed-style runs. Each run is (text, font_size, bold, color)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    for i, (text, font_size, bold, color) in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = 'Calibri'
    return txBox


def add_rect(slide, left, top, width, height, fill_color):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_slide_number(slide, num, color=GREY):
    """Add slide number bottom-right."""
    add_textbox(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.55),
                Inches(0.8), Inches(0.3), f"{num:02d}",
                font_size=10, color=color, alignment=PP_ALIGN.RIGHT)


def add_logo(slide, variant='black'):
    """Add footer logo. variant='black' for light bg, 'white' for dark bg."""
    logo_path = os.path.join(os.path.dirname(__file__), f"logo-footer-{variant}.png")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.7), SLIDE_H - Inches(0.6),
                                 width=Inches(1.8))


def set_bg(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, WHITE)

add_textbox(slide, MARGIN, Inches(2.0), Inches(6), Inches(0.3),
            "TEAM TALK  /  FEBRUARY 2026", font_size=10, bold=True, color=GREY)

add_textbox(slide, MARGIN, Inches(2.8), Inches(10), Inches(1.2),
            "Where We Stand", font_size=60, bold=False, color=BLACK)

add_textbox(slide, MARGIN, Inches(4.2), Inches(6), Inches(0.8),
            "A hard look at our position. No sugarcoating. No panic. Clear-eyed assessment from partner to partner.",
            font_size=18, color=GREY, line_spacing=28)

add_logo(slide, 'black')
add_slide_number(slide, 1)


# ============================================================
# SLIDE 2: THE NUMBERS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_LIGHT)

add_textbox(slide, MARGIN, Inches(0.7), Inches(6), Inches(0.3),
            "THE HONEST NUMBERS", font_size=10, bold=True, color=GREY)

add_textbox(slide, MARGIN, Inches(1.1), Inches(10), Inches(0.6),
            "Where our funds sit today", font_size=28, bold=False, color=BLACK)

# Fund cards - 3 columns
card_w = Inches(3.7)
card_h = Inches(4.5)
card_y = Inches(2.0)
gap = Inches(0.3)

# Fund I
x = MARGIN
add_rect(slide, x, card_y, card_w, card_h, WHITE)
add_textbox(slide, x + Inches(0.3), card_y + Inches(0.3), card_w - Inches(0.6), Inches(0.25),
            "FUND I  /  PILOT", font_size=9, bold=True, color=RGBColor(0x8A, 0x9A, 0x8D))
add_textbox(slide, x + Inches(0.3), card_y + Inches(0.7), Inches(2), Inches(0.7),
            "8x", font_size=48, bold=True, color=RGBColor(0x8A, 0x9A, 0x8D))
add_textbox(slide, x + Inches(0.3), card_y + Inches(1.4), Inches(2), Inches(0.3),
            "TVPI", font_size=14, color=GREY)
add_textbox(slide, x + Inches(0.3), card_y + Inches(2.0), Inches(1.2), Inches(0.2),
            "DPI", font_size=10, color=GREY)
add_textbox(slide, x + Inches(0.3), card_y + Inches(2.2), Inches(1.2), Inches(0.4),
            "3x", font_size=26, bold=True, color=BLACK)
add_textbox(slide, x + Inches(1.8), card_y + Inches(2.0), Inches(1.5), Inches(0.2),
            "Fund Size", font_size=10, color=GREY)
add_textbox(slide, x + Inches(1.8), card_y + Inches(2.2), Inches(1.5), Inches(0.4),
            "$1.1M", font_size=26, bold=True, color=BLACK)
add_textbox(slide, x + Inches(0.3), card_y + Inches(3.0), card_w - Inches(0.6), Inches(1.2),
            "Exceptional returns. Real distributions. But LPs discount the fund size as \"angel-scale, not institutional.\"",
            font_size=13, color=GREY, line_spacing=20)

# Fund II
x = MARGIN + card_w + gap
add_rect(slide, x, card_y, card_w, card_h, WHITE)
add_textbox(slide, x + Inches(0.3), card_y + Inches(0.3), card_w - Inches(0.6), Inches(0.25),
            "FUND II  /  MERCURY", font_size=9, bold=True, color=GREY)
add_textbox(slide, x + Inches(0.3), card_y + Inches(0.7), Inches(2.5), Inches(0.7),
            "0.95x", font_size=48, bold=True, color=GREY)
add_textbox(slide, x + Inches(0.3), card_y + Inches(1.4), Inches(2), Inches(0.3),
            "TVPI (current)", font_size=14, color=GREY)
add_textbox(slide, x + Inches(0.3), card_y + Inches(2.0), Inches(1.5), Inches(0.2),
            "Bullish scenario", font_size=10, color=GREY)
add_textbox(slide, x + Inches(0.3), card_y + Inches(2.2), Inches(1.5), Inches(0.4),
            "~1.3x", font_size=26, bold=True, color=BLACK)
add_textbox(slide, x + Inches(2.0), card_y + Inches(2.0), Inches(1.5), Inches(0.2),
            "Fund Size", font_size=10, color=GREY)
add_textbox(slide, x + Inches(2.0), card_y + Inches(2.2), Inches(1.5), Inches(0.4),
            "$40M", font_size=26, bold=True, color=BLACK)
add_textbox(slide, x + Inches(0.3), card_y + Inches(3.0), card_w - Inches(0.6), Inches(1.2),
            "2022 vintage. Pre-seed/seed deep tech. Below water on marks. This is the number LPs see first.",
            font_size=13, color=GREY, line_spacing=20)

# Fund III
x = MARGIN + 2 * (card_w + gap)
add_rect(slide, x, card_y, card_w, card_h, CARD_AZURE)
add_textbox(slide, x + Inches(0.3), card_y + Inches(0.3), card_w - Inches(0.6), Inches(0.25),
            "FUND III  /  MARS", font_size=9, bold=True, color=AZURE_DARK)
add_textbox(slide, x + Inches(0.3), card_y + Inches(0.7), Inches(2), Inches(0.7),
            "$5M", font_size=48, bold=True, color=AZURE_DARK)
add_textbox(slide, x + Inches(0.3), card_y + Inches(1.4), Inches(2), Inches(0.3),
            "First Close", font_size=14, color=GREY)
add_textbox(slide, x + Inches(0.3), card_y + Inches(2.0), Inches(1.2), Inches(0.2),
            "Target", font_size=10, color=GREY)
add_textbox(slide, x + Inches(0.3), card_y + Inches(2.2), Inches(1.2), Inches(0.4),
            "$50M", font_size=26, bold=True, color=BLACK)
add_textbox(slide, x + Inches(1.8), card_y + Inches(2.0), Inches(1.5), Inches(0.2),
            "Soft Commits", font_size=10, color=GREY)
add_textbox(slide, x + Inches(1.8), card_y + Inches(2.2), Inches(1.5), Inches(0.4),
            "~$2M", font_size=26, bold=True, color=BLACK)
add_textbox(slide, x + Inches(0.3), card_y + Inches(3.0), card_w - Inches(0.6), Inches(1.2),
            "Long way to go. 16 data room LPs ghosted. The fundraise is stalling.",
            font_size=13, color=GREY, line_spacing=20)

add_logo(slide, 'black')
add_slide_number(slide, 2)


# ============================================================
# SLIDE 3: CONTEXT
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, WHITE)

add_textbox(slide, MARGIN, Inches(0.7), Inches(6), Inches(0.3),
            "CONTEXT", font_size=10, bold=True, color=GREY)

add_textbox(slide, MARGIN, Inches(1.1), Inches(10), Inches(0.8),
            "1.3x TVPI at year 4 is not a death sentence.\nBut it is not enough to raise on numbers alone.",
            font_size=28, color=BLACK, line_spacing=38)

# Left column - 3 grey cards
col_w = Inches(5.5)
y_start = Inches(2.5)
card_h_sm = Inches(1.3)

items_left = [
    ("2022 is the worst vintage in a decade",
     "Valuations peaked, markets crashed, follow-on funding dried up. Every 2022 vintage fund is sitting on compressed marks. This is not unique to us."),
    ("Deep tech J-curve is longer",
     "SaaS marks up on revenue multiples. Deep tech marks up on technical milestones, government contracts, and large follow-on rounds. Our companies are pre-inflection, not failing."),
    ("Funds settle into terminal quartile at year 5-6",
     "Cambridge Associates data is clear: judging a deep tech fund at year 3-4 is like judging a marathon runner at mile 8."),
]

for i, (title, body) in enumerate(items_left):
    y = y_start + i * (card_h_sm + Inches(0.2))
    add_rect(slide, MARGIN, y, col_w, card_h_sm, BG_GREY)
    add_textbox(slide, MARGIN + Inches(0.3), y + Inches(0.15), col_w - Inches(0.6), Inches(0.3),
                title, font_size=15, bold=True, color=BLACK)
    add_textbox(slide, MARGIN + Inches(0.3), y + Inches(0.5), col_w - Inches(0.6), Inches(0.7),
                body, font_size=13, color=GREY, line_spacing=20)

# Right column - key insight card
rx = MARGIN + col_w + Inches(0.4)
rw = Inches(5.5)

add_rect(slide, rx, y_start, rw, Inches(1.3), WHITE)
add_textbox(slide, rx + Inches(0.3), y_start + Inches(0.15), rw - Inches(0.6), Inches(0.3),
            "But median does not raise a Fund III", font_size=15, bold=True, color=BLACK)
add_textbox(slide, rx + Inches(0.3), y_start + Inches(0.5), rw - Inches(0.6), Inches(0.7),
            "1.3x at year 4 is roughly median for the vintage. LPs want top quartile trajectory or a compelling story for why marks are about to inflect.",
            font_size=13, color=GREY, line_spacing=20)

# Verdant highlight card
vy = y_start + Inches(1.6)
add_rect(slide, rx, vy, rw, Inches(2.7), CARD_VERDANT)
add_textbox(slide, rx + Inches(0.3), vy + Inches(0.2), rw - Inches(0.6), Inches(0.3),
            "What actually matters", font_size=15, bold=True, color=VERDANT_DARK)
add_textbox(slide, rx + Inches(0.3), vy + Inches(0.6), rw - Inches(0.6), Inches(1.0),
            "It is not the absolute number. It is the slope of the curve and the credibility of what comes next.",
            font_size=17, bold=False, color=BLACK, line_spacing=26)
add_textbox(slide, rx + Inches(0.3), vy + Inches(1.9), rw - Inches(0.6), Inches(0.7),
            "If we show LPs that 0.95x is becoming 1.3x because specific companies received term sheets at validated external prices, that momentum signal is more persuasive than the number itself.",
            font_size=12, color=GREY, line_spacing=18)

add_logo(slide, 'black')
add_slide_number(slide, 3)


# ============================================================
# SLIDE 4: DIAGNOSIS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_LIGHT)

add_textbox(slide, MARGIN, Inches(0.7), Inches(6), Inches(0.3),
            "DIAGNOSIS", font_size=10, bold=True, color=GREY)
add_textbox(slide, MARGIN, Inches(1.1), Inches(10), Inches(0.6),
            "Why 16 out of 16 data room LPs went silent", font_size=28, color=BLACK)

items = [
    ("1", "Fund II gives them permission to wait",
     "LP opens data room, sees below-water marks, thinks: \"I will wait for momentum.\" Ghosting is easier than saying \"your numbers are not there yet.\""),
    ("2", "No urgency mechanism in our materials",
     "Our pitch reads as standard deep tech thesis + team + portfolio. Nothing makes LPs feel they will miss something irreplaceable by waiting."),
    ("3", "Allocation bucket mismatch",
     "\"European pre-seed deep tech\" is not a standard allocation category. We are asking LPs to create a new bucket. That is a fundamentally harder sale."),
]

for i, (num, title, body) in enumerate(items):
    y = Inches(2.2) + i * Inches(1.6)
    add_textbox(slide, MARGIN, y, Inches(0.5), Inches(0.5),
                num, font_size=32, bold=True, color=GREY)
    add_textbox(slide, MARGIN + Inches(0.6), y, Inches(6), Inches(0.35),
                title, font_size=17, bold=True, color=BLACK)
    add_textbox(slide, MARGIN + Inches(0.6), y + Inches(0.4), Inches(6), Inches(0.8),
                body, font_size=13, color=GREY, line_spacing=20)

# Compounding effect card
cx = Inches(8.2)
cw = Inches(4.3)
add_rect(slide, cx, Inches(2.5), cw, Inches(3.5), CARD_AZURE)
add_textbox(slide, cx + Inches(0.3), Inches(2.7), cw - Inches(0.6), Inches(0.3),
            "The compounding effect", font_size=15, bold=True, color=BLACK)
add_textbox(slide, cx + Inches(0.3), Inches(3.1), cw - Inches(0.6), Inches(1.5),
            "No single factor is a dealbreaker. But combined, they create inertia. Performance gives permission to wait. Lack of differentiation gives no reason to act. Bucket mismatch removes the internal champion.",
            font_size=13, color=GREY, line_spacing=20)
add_textbox(slide, cx + Inches(0.3), Inches(4.8), cw - Inches(0.6), Inches(0.6),
            "We need to break the compound. Attack all three simultaneously.",
            font_size=13, bold=True, color=BLACK, line_spacing=20)

add_logo(slide, 'black')
add_slide_number(slide, 4)


# ============================================================
# SLIDE 5: BEING HONEST (split)
# ============================================================
slide = prs.slides.add_slide(blank_layout)

# Left black panel
half_w = Inches(6.667)
add_rect(slide, Inches(0), Inches(0), half_w, SLIDE_H, BLACK)

add_textbox(slide, Inches(0.7), Inches(2.0), Inches(5.5), Inches(0.3),
            "BEING HONEST WITH YOU", font_size=10, bold=True, color=GREY)
add_textbox(slide, Inches(0.7), Inches(2.6), Inches(5.5), Inches(1.2),
            "I am not sleeping well.\nI am scared the team falls apart.\nI owe you honesty.",
            font_size=26, color=WHITE, line_spacing=38)
add_textbox(slide, Inches(0.7), Inches(4.2), Inches(5.2), Inches(2.0),
            "The org is fully capitalised until mid-2027. From there on we have to cut costs. This does not take into account management fees from Mars. The pressure is real. I am telling you this because I believe that transparency and agency are our super powers.",
            font_size=15, color=GREY, line_spacing=24)

# Right white panel
rx = half_w
set_bg(slide, WHITE)
# Need to re-draw black since bg is white
add_rect(slide, Inches(0), Inches(0), half_w, SLIDE_H, BLACK)

add_textbox(slide, rx + Inches(0.7), Inches(1.8), Inches(5.5), Inches(0.3),
            "WHAT I HAVE BEEN WRESTLING WITH", font_size=10, bold=True, color=GREY)

items_r = [
    ("The existential conflation",
     "When your life's work feels threatened, \"this fundraise is hard\" becomes \"everything might fail.\" These are different problems. I am working to separate them."),
    ("The temptation to scatter",
     "Fear pushes toward consensus solutions - launch products, do consulting, hedge bets. That is the playbook investing we left to escape."),
    ("The thing I know to be true",
     "We built something genuinely different. Fund I at 8x proved the model. The portfolio companies are real. The thesis is playing out. The packaging is wrong, not the substance."),
]

for i, (title, body) in enumerate(items_r):
    y = Inches(2.4) + i * Inches(1.5)
    tc = VERDANT_DARK if i == 2 else BLACK
    add_textbox(slide, rx + Inches(0.7), y, Inches(5.5), Inches(0.3),
                title, font_size=15, bold=True, color=tc)
    add_textbox(slide, rx + Inches(0.7), y + Inches(0.35), Inches(5.2), Inches(0.8),
                body, font_size=13, color=GREY, line_spacing=20)

add_logo(slide, 'white')
add_slide_number(slide, 5)


# ============================================================
# SLIDE 6: ON AGENCY
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, WHITE)

add_textbox(slide, MARGIN, Inches(0.9), Inches(6), Inches(0.3),
            "ON AGENCY", font_size=10, bold=True, color=GREY)
add_textbox(slide, MARGIN, Inches(1.5), Inches(10), Inches(0.9),
            "We can turn things around.\nWe thrive in adversity.",
            font_size=40, color=BLACK, line_spacing=52)

# Quote card
qy = Inches(3.0)
qw = Inches(11.5)
qh = Inches(3.8)
add_rect(slide, MARGIN, qy, qw, qh, CARD_VERDANT)

add_textbox(slide, MARGIN + Inches(0.4), qy + Inches(0.3), qw - Inches(0.8), Inches(0.3),
            "FROM OUR FOUNDER JOB", font_size=9, bold=True, color=VERDANT_DARK)

quote1 = "\"Monday is the best day of the week. It's our chance to start fresh, to improve upon who we were last week. Today is a tougher Monday than usual: obviously, the situation with Brandon is challenging. But the most important question right now is this: Who are we going to be in the face of adversity?"
quote2 = "For Hedy to fulfill its mission and achieve its vision, there's only one answer: we love adversity. We thrive in it. Adversity makes us stronger, sharper, and better. Like iron in a blast furnace, adversity burns away the weak points and forges steel."
quote3 = "The status quo isn't good enough today. Who we were last week won't cut it anymore. Activate your inner DARKMODE. Get better this week.\""

add_textbox(slide, MARGIN + Inches(0.4), qy + Inches(0.7), qw - Inches(0.8), Inches(1.0),
            quote1, font_size=14, color=BLACK, line_spacing=22)
add_textbox(slide, MARGIN + Inches(0.4), qy + Inches(1.7), qw - Inches(0.8), Inches(0.9),
            quote2, font_size=14, color=BLACK, line_spacing=22)
add_textbox(slide, MARGIN + Inches(0.4), qy + Inches(2.7), qw - Inches(0.8), Inches(0.9),
            quote3, font_size=14, color=BLACK, line_spacing=22)

add_logo(slide, 'black')
add_slide_number(slide, 6)


# ============================================================
# SLIDE 7: A REQUEST
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_LIGHT)

add_textbox(slide, MARGIN, Inches(0.7), Inches(6), Inches(0.3),
            "A REQUEST", font_size=10, bold=True, color=GREY)
add_textbox(slide, MARGIN, Inches(1.2), Inches(10), Inches(0.9),
            "I am not asking you to stay out of loyalty.\nI am asking you to stay out of conviction.",
            font_size=38, color=BLACK, line_spacing=50)

# Two cards
cw7 = Inches(5.7)
cy7 = Inches(2.8)
ch7 = Inches(3.8)

# Card 1: Catalysts
add_rect(slide, MARGIN, cy7, cw7, ch7, WHITE)
add_textbox(slide, MARGIN + Inches(0.3), cy7 + Inches(0.25), cw7 - Inches(0.6), Inches(0.35),
            "Help manufacture catalysts", font_size=17, bold=True, color=BLACK)

# Mercury paragraph
txBox = slide.shapes.add_textbox(MARGIN + Inches(0.3), cy7 + Inches(0.75), cw7 - Inches(0.6), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.line_spacing = Pt(20)
run_b = p.add_run()
run_b.text = "Mercury: "
run_b.font.size = Pt(13)
run_b.font.bold = True
run_b.font.color.rgb = BLACK
run_b.font.name = 'Calibri'
run_t = p.add_run()
run_t.text = "We need 2-3 portfolio companies to hit visible milestones in the next 90 days. Term sheets, contracts, partnerships. This changes the fundraising narrative from \"wait and see\" to \"momentum is building.\" Your relationships and support matter here."
run_t.font.size = Pt(13)
run_t.font.color.rgb = GREY
run_t.font.name = 'Calibri'

# Mars paragraph
txBox2 = slide.shapes.add_textbox(MARGIN + Inches(0.3), cy7 + Inches(2.3), cw7 - Inches(0.6), Inches(0.8))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.line_spacing = Pt(20)
run_b2 = p2.add_run()
run_b2.text = "Mars: "
run_b2.font.size = Pt(13)
run_b2.font.bold = True
run_b2.font.color.rgb = BLACK
run_b2.font.name = 'Calibri'
run_t2 = p2.add_run()
run_t2.text = "We need to do 1-2 bold earliest stage investments at $1M each to build momentum with the Mars Fund."
run_t2.font.size = Pt(13)
run_t2.font.color.rgb = GREY
run_t2.font.name = 'Calibri'

# Card 2: Energy
x2 = MARGIN + cw7 + Inches(0.3)
add_rect(slide, x2, cy7, cw7, ch7, WHITE)
add_textbox(slide, x2 + Inches(0.3), cy7 + Inches(0.25), cw7 - Inches(0.6), Inches(0.35),
            "Protect your energy", font_size=17, bold=True, color=BLACK)
add_textbox(slide, x2 + Inches(0.3), cy7 + Inches(0.75), cw7 - Inches(0.6), Inches(2.5),
            "We are a small team carrying a big mission. The next 12 months will not be easy. No side projects. No non-portfolio SPVs. No hedging through alternative revenues. No distractions. We give this a 100% shot.",
            font_size=13, color=GREY, line_spacing=20)

add_logo(slide, 'black')
add_slide_number(slide, 7)


# ============================================================
# SLIDE 8: THE PLAN
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, WHITE)

add_textbox(slide, MARGIN, Inches(0.6), Inches(6), Inches(0.3),
            "THE PATH FORWARD", font_size=10, bold=True, color=GREY)
add_textbox(slide, MARGIN, Inches(1.0), Inches(10), Inches(0.5),
            "Seven moves. In order. No shortcuts.", font_size=28, color=BLACK)

steps = [
    ("1", "Fix the data room", "February. Get brutal feedback from Hutt and Isomer. What exactly triggers the \"wait\"?", "NOW"),
    ("2", "Lead with identity, not data", "Feb-March. Rewrite the pitch. Who we are first. Numbers validate the story.", "FEB"),
    ("3", "Manufacture portfolio catalysts", "March-Oct. Pick 2-3 Mercury companies closest to milestones. Focus all support there for 90 days. Do 2-3 Mars Fund investments of $1M each.", "MAR"),
    ("4", "Launch follow-on SPVs", "Q2-Q3. When catalysts hit, offer LPs co-investment into our own companies. Not random deals.", "Q2"),
    ("5", "Kill the side projects", "Now. No marketplace. No consulting. Every hour goes to fundraising, portfolio, and deals.", "NOW"),
    ("6", "Re-engage ghosted LPs with momentum", "May. Push updated marks + new pitch to all 16 data room LPs. Different story, different energy.", "MAY"),
    ("7", "Find peer GPs for mutual support", "Ongoing. 2-3 emerging managers in similar position. Isolation amplifies fear. Peer networks amplify clarity.", "ONGOING"),
]

col_w = Inches(5.7)
row_h = Inches(0.88)

for i, (num, title, desc, timing) in enumerate(steps):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    x = MARGIN + col * (col_w + Inches(0.3))
    y = Inches(1.8) + row * (row_h + Inches(0.08))

    add_rect(slide, x, y, col_w, row_h, BG_GREY)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.08), Inches(0.4), Inches(0.4),
                num, font_size=22, bold=True, color=BLACK)
    add_textbox(slide, x + Inches(0.55), y + Inches(0.06), col_w - Inches(1.5), Inches(0.3),
                title, font_size=14, bold=True, color=BLACK)
    add_textbox(slide, x + Inches(0.55), y + Inches(0.34), col_w - Inches(1.5), Inches(0.52),
                desc, font_size=10, color=GREY, line_spacing=15)
    add_textbox(slide, x + col_w - Inches(1.0), y + Inches(0.28), Inches(0.8), Inches(0.25),
                timing, font_size=8, bold=True, color=GREY, alignment=PP_ALIGN.RIGHT)

# Target outcome card
ty = Inches(5.7)
add_rect(slide, MARGIN, ty, col_w, Inches(0.95), CARD_VERDANT)
add_textbox(slide, MARGIN + Inches(0.3), ty + Inches(0.08), col_w - Inches(0.6), Inches(0.25),
            "Target outcome", font_size=12, bold=True, color=VERDANT_DARK)
add_textbox(slide, MARGIN + Inches(0.3), ty + Inches(0.35), col_w - Inches(0.6), Inches(0.55),
            "Second close at $15-20M by summer. Third close at $25-35M by year-end. Momentum narrative replaces \"wait and see.\" To hit $50M+ we can extend the fundraising timeline into 2027 (subject to LPA).",
            font_size=10, color=GREY, line_spacing=15)

add_logo(slide, 'black')
add_slide_number(slide, 8)


# ============================================================
# SLIDE 9: TRANSPARENCY
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_LIGHT)

add_textbox(slide, MARGIN, Inches(0.7), Inches(6), Inches(0.3),
            "TRANSPARENCY", font_size=10, bold=True, color=GREY)
add_textbox(slide, MARGIN, Inches(1.1), Inches(10), Inches(0.6),
            "Ideas we considered. And why we should question them.", font_size=28, color=BLACK)

cards_t = [
    ("Defense procurement marketplace", "DISCUSS",
     "Building a marketplace is a full-time job for a dedicated team. Creates conflict of interest with portfolio companies. If it fails, we burned 12-18 months of bandwidth. If it works, LPs question whether we are a fund or an operator.",
     WHITE, GREY),
    ("AI consulting for Mittelstand", "DISCUSS",
     "Consulting trades time for money at a linear rate. Zero convexity. Every hour consulting is an hour not fundraising. Signals to LPs that the fund cannot sustain itself. But it does generate short-term cash flow.",
     WHITE, GREY),
    ("SPVs into elite companies", "DISCUSS",
     "SPVs into companies we have no relationship with is a different skill set. Dilutes our brand. Hard to get allocation. But follow-on SPVs into our own portfolio could validate our picks, give LPs a low-risk entry, and extend runway.",
     CARD_AZURE, AZURE_DARK),
]

cw_t = Inches(3.7)
for i, (title, label, body, bg, label_color) in enumerate(cards_t):
    x = MARGIN + i * (cw_t + Inches(0.3))
    y = Inches(2.2)
    add_rect(slide, x, y, cw_t, Inches(4.5), bg)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), cw_t - Inches(0.6), Inches(0.35),
                title, font_size=17, bold=True, color=BLACK)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.75), cw_t - Inches(0.6), Inches(0.25),
                label, font_size=9, bold=True, color=label_color)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.2), cw_t - Inches(0.6), Inches(2.8),
                body, font_size=13, color=GREY, line_spacing=20)

add_logo(slide, 'black')
add_slide_number(slide, 9)


# ============================================================
# SLIDE 10: HONEY BADGER MINDSET
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, VERDANT)

add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(0.3),
            "HONEY BADGER MINDSET", font_size=10, bold=True, color=BLACK,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.0),
            "Fear whispers:\n\"Pivot to safety. Diversify. Hedge your bets.\"",
            font_size=32, color=BLACK, alignment=PP_ALIGN.CENTER, line_spacing=44)

add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(1.4),
            "That is what corporates would do.\nWe all left them to escape and build\nsomething meaningful.",
            font_size=32, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER, line_spacing=44)

add_textbox(slide, Inches(2.0), Inches(5.2), Inches(9), Inches(1.2),
            "The honey badger move now is to attack. Double down on who we are, fix the packaging, manufacture the catalysts, and let the market validate what we already know about our portfolio.",
            font_size=16, color=BLACK, alignment=PP_ALIGN.CENTER, line_spacing=26)

add_logo(slide, 'black')
add_slide_number(slide, 10, VERDANT_DARK)


# ============================================================
# SLIDE 11: CLOSING (split)
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, WHITE)

half_w = Inches(6.667)
add_rect(slide, Inches(0), Inches(0), half_w, SLIDE_H, BLACK)

add_textbox(slide, Inches(0.7), Inches(2.2), Inches(5.5), Inches(0.3),
            "CLOSING", font_size=10, bold=True, color=GREY)
add_textbox(slide, Inches(0.7), Inches(2.8), Inches(5.5), Inches(1.0),
            "This is not a crisis.\nThis is a test.",
            font_size=38, color=WHITE, line_spacing=50)
add_textbox(slide, Inches(0.7), Inches(4.2), Inches(5.2), Inches(1.5),
            "The hardest transition in fund management is micro to institutional. The mortality rate is high. But so is our conviction. We have done hard things before. We will do this one together.",
            font_size=15, color=GREY, line_spacing=24)

# Right panel items
items_close = [
    ("We are not dying",
     "We have runway, committed LPs, a proven thesis, and a portfolio that is pre-inflection, not failing."),
    ("The substance is right",
     "8x Fund I. Sovereign compute thesis playing out. Defense pivot ahead of market. The packaging is wrong, not the product."),
    ("We go forward together",
     "No information asymmetry. Same numbers, same pressure, same mission. For us this is a life's work. That has not changed."),
]

for i, (title, body) in enumerate(items_close):
    y = Inches(2.0) + i * Inches(1.4)
    add_textbox(slide, half_w + Inches(0.7), y, Inches(5.5), Inches(0.3),
                title, font_size=17, bold=True, color=BLACK)
    add_textbox(slide, half_w + Inches(0.7), y + Inches(0.35), Inches(5.2), Inches(0.7),
                body, font_size=13, color=GREY, line_spacing=20)

add_textbox(slide, half_w + Inches(0.7), Inches(6.2), Inches(5), Inches(0.4),
            "inflection", font_size=28, bold=True, color=BLACK)
add_textbox(slide, half_w + Inches(0.7), Inches(6.6), Inches(3), Inches(0.3),
            "February 2026", font_size=13, color=GREY)

add_logo(slide, 'white')
add_slide_number(slide, 11)


# Save
output_path = os.path.join(os.path.dirname(__file__), "mars-team-talk.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
