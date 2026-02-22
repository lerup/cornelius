#!/usr/bin/env python3
"""Export HTML pitch deck to PDF and editable PPTX."""

import os, time, tempfile
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DECK_DIR = Path(__file__).parent
HTML_PATH = DECK_DIR / "section-0-1.html"
OUT_NAME = "Mars Pitch Hard Futures 20260214"
PDF_PATH = DECK_DIR / f"{OUT_NAME}.pdf"
PPTX_PATH = DECK_DIR / f"{OUT_NAME}.pptx"

# 16:9 dimensions
SLIDE_W = 1920
SLIDE_H = 1080


def export_pdf():
    """Use Playwright to render each slide and combine into a single PDF."""
    print("Generating PDF...")
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": SLIDE_W, "height": SLIDE_H})
        page.goto(f"file://{HTML_PATH.resolve()}")
        page.wait_for_load_state("networkidle")
        time.sleep(1)

        # Get total slide count
        total = page.evaluate("document.querySelectorAll('.slide').length")
        print(f"  Found {total} slides")

        # Screenshot each slide into a temp dir
        tmp = tempfile.mkdtemp()
        paths = []
        for i in range(total):
            page.evaluate(f"""
                document.querySelectorAll('.slide').forEach((s, idx) => {{
                    s.classList.toggle('active', idx === {i});
                    s.style.display = idx === {i} ? 'flex' : 'none';
                }});
                // Force deck to fill viewport
                document.querySelector('.deck').style.width = '{SLIDE_W}px';
                document.querySelector('.deck').style.height = '{SLIDE_H}px';
                // Hide nav elements
                document.querySelectorAll('.nav-hint, .slide-counter').forEach(el => el.style.display = 'none');
            """)
            time.sleep(0.3)
            path = os.path.join(tmp, f"slide_{i:03d}.png")
            # Screenshot the deck element
            deck = page.query_selector(".deck")
            deck.screenshot(path=path)
            paths.append(path)
            print(f"  Captured slide {i}")

        browser.close()

    # Combine PNGs into PDF using img2pdf or reportlab
    try:
        import img2pdf
        with open(PDF_PATH, "wb") as f:
            f.write(img2pdf.convert(paths))
    except ImportError:
        # Fallback: use Pillow
        from PIL import Image
        images = [Image.open(p).convert("RGB") for p in paths]
        images[0].save(PDF_PATH, save_all=True, append_images=images[1:], resolution=150)

    # Cleanup
    for p in paths:
        os.remove(p)
    os.rmdir(tmp)
    print(f"  PDF saved: {PDF_PATH}")
    return paths


def export_pptx(slide_images=None):
    """Create editable PPTX by placing slide screenshots as backgrounds
    with editable text overlays for key content."""
    print("Generating PPTX...")

    # First, capture fresh screenshots if not provided
    tmp = tempfile.mkdtemp()
    img_paths = []

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": SLIDE_W, "height": SLIDE_H})
        page.goto(f"file://{HTML_PATH.resolve()}")
        page.wait_for_load_state("networkidle")
        time.sleep(1)

        total = page.evaluate("document.querySelectorAll('.slide').length")

        for i in range(total):
            page.evaluate(f"""
                document.querySelectorAll('.slide').forEach((s, idx) => {{
                    s.classList.toggle('active', idx === {i});
                    s.style.display = idx === {i} ? 'flex' : 'none';
                }});
                document.querySelector('.deck').style.width = '{SLIDE_W}px';
                document.querySelector('.deck').style.height = '{SLIDE_H}px';
                document.querySelectorAll('.nav-hint, .slide-counter').forEach(el => el.style.display = 'none');
            """)
            time.sleep(0.3)
            path = os.path.join(tmp, f"pptx_slide_{i:03d}.png")
            deck = page.query_selector(".deck")
            deck.screenshot(path=path)
            img_paths.append(path)

        # Also extract text content from each slide for editable overlays
        slide_texts = []
        for i in range(total):
            page.evaluate(f"""
                document.querySelectorAll('.slide').forEach((s, idx) => {{
                    s.classList.toggle('active', idx === {i});
                    s.style.display = idx === {i} ? 'flex' : 'none';
                }});
            """)
            h2 = page.evaluate(f"document.querySelectorAll('.slide')[{i}].querySelector('h2')?.innerText || ''")
            slide_texts.append({"title": h2})

        browser.close()

    # Create PPTX
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W * 9525)   # EMU = 1/914400 inch, 9525 EMU per pixel at 96dpi
    prs.slide_height = Emu(SLIDE_H * 9525)

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    for i, img_path in enumerate(img_paths):
        slide = prs.slides.add_slide(blank_layout)

        # Add screenshot as full-bleed background image
        slide.shapes.add_picture(
            img_path,
            Emu(0), Emu(0),
            prs.slide_width, prs.slide_height
        )

        # Add invisible editable text box with slide title (for search/edit)
        title_text = slide_texts[i]["title"] if i < len(slide_texts) else ""
        if title_text:
            txBox = slide.shapes.add_textbox(
                Emu(int(72 * 9525)),   # left padding matching CSS
                Emu(int(48 * 9525)),   # top padding matching CSS
                Emu(int(1600 * 9525)), # width
                Emu(int(80 * 9525))    # height
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = title_text
            run.font.size = Pt(1)  # Tiny - just for editability/search
            run.font.color.rgb = RGBColor(255, 255, 255)  # White = invisible on white
            # Make textbox transparent
            txBox.fill.background()

        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = f"Slide {i}"

    prs.save(str(PPTX_PATH))

    # Cleanup
    for p in img_paths:
        os.remove(p)
    os.rmdir(tmp)
    print(f"  PPTX saved: {PPTX_PATH}")


if __name__ == "__main__":
    export_pdf()
    export_pptx()
    print(f"\nDone. Files in: {DECK_DIR}")
