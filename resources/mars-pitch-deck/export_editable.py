#!/usr/bin/env python3
"""Export HTML pitch deck to fully editable PPTX with native elements."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import copy

DECK_DIR = Path(__file__).parent
PPTX_PATH = DECK_DIR / "Mars Pitch Hard Futures 20260213.pptx"

# 16:9 in EMUs (standard widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Brand colors
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(0xB3, 0xBC, 0xB5)
BLUE = RGBColor(0xB4, 0xBA, 0xCC)
GREY = RGBColor(0x80, 0x80, 0x80)
LIGHT_GREY = RGBColor(0xAA, 0xAA, 0xAA)
DARK_GREY = RGBColor(0x55, 0x55, 0x55)
BG_GREY = RGBColor(0xF0, 0xF0, 0xF0)
BORDER_GREY = RGBColor(0xE0, 0xE0, 0xDE)
DARK_GREEN = RGBColor(0x6A, 0x8A, 0x6E)

# Font
FONT = "Inter"

# Padding
PAD_L = Inches(0.75)
PAD_T = Inches(0.5)
PAD_R = Inches(0.75)
CONTENT_W = SLIDE_W - PAD_L - PAD_R


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False,
                 color=BLACK, alignment=PP_ALIGN.LEFT, font_name=FONT, line_spacing=1.15):
    """Helper to add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multi_text(slide, left, top, width, height, segments, font_size=14,
                   alignment=PP_ALIGN.LEFT, line_spacing=1.15):
    """Add textbox with multiple styled runs. segments = [(text, bold, color), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(0)
    for text, bold, color in segments:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
    return txBox


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_logo(slide, color="black"):
    """Add logo text placeholder top-left."""
    c = GREY if color == "black" else RGBColor(0xCC, 0xCC, 0xCC)
    add_text_box(slide, PAD_L, Inches(0.4), Inches(2), Inches(0.5),
                 "inflection", font_size=20, bold=True, color=c)


def slide_title(slide, text, top=Inches(1.0), color=BLACK, size=32):
    """Add slide title."""
    add_text_box(slide, PAD_L, top, CONTENT_W, Inches(0.8), text,
                 font_size=size, bold=True, color=color)


# ═══════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════

def slide_0_title(prs):
    """Title slide - black bg."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BLACK

    add_text_box(slide, Inches(0), Inches(1.5), SLIDE_W, Inches(1),
                 "inflection", font_size=48, bold=True, color=RGBColor(0xCC, 0xCC, 0xCC),
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(3.0), SLIDE_W, Inches(1.2),
                 "Hard Futures", font_size=72, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(4.3), SLIDE_W, Inches(0.5),
                 "FUND III  -  MARS", font_size=18, color=GREY,
                 alignment=PP_ALIGN.CENTER)


def slide_1_origins(prs):
    """Origins of venture - dark bg with text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x11, 0x11, 0x11)

    # Try to add the Fairchild image
    img_path = DECK_DIR / "origins-fairchild.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Emu(0), Emu(0), SLIDE_W, SLIDE_H)
        # Dark overlay (solid black, opacity via XML)
        overlay = add_shape(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill_color=BLACK)
        from pptx.oxml.ns import qn
        from lxml import etree
        sp_pr = overlay._element.find(qn('a:solidFill'))
        # Navigate shape XML to set alpha
        for sf in overlay._element.iter(qn('a:solidFill')):
            for clr in sf:
                alpha = etree.SubElement(clr, qn('a:alpha'))
                alpha.set('val', '70000')
                break

    add_text_box(slide, PAD_L, Inches(4.0), Inches(8), Inches(0.8),
                 "The origins of venture", font_size=40, bold=True, color=WHITE)

    add_multi_text(slide, PAD_L, Inches(4.8), Inches(9), Inches(0.6), [
        ("1957", True, GREEN),
        (" - Arthur Rock backed eight scientists who walked out of Shockley Semiconductor. No business plan. No pitch deck. $1.38M from Fairchild Camera. That bet created Intel, AMD, and the entire semiconductor industry.", False, RGBColor(0xCC,0xCC,0xCC)),
    ], font_size=13)

    add_multi_text(slide, PAD_L, Inches(5.5), Inches(9), Inches(0.5), [
        ("1976", True, GREEN),
        (" - Mike Markkula drove to a garage and wrote a $250,000 check for a third of Apple. Two kids and a circuit board. No IC committee. No consensus.", False, RGBColor(0xCC,0xCC,0xCC)),
    ], font_size=13)

    add_text_box(slide, PAD_L, Inches(6.2), Inches(10), Inches(0.5),
                 "Every legendary venture return came from the same place: one person's conviction about something most people could not yet evaluate.",
                 font_size=14, bold=False, color=RGBColor(0x99, 0x99, 0x99))


def slide_2_concentration(prs):
    """PE playbooks + concentration data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    slide_title(slide, "Today the industry runs\non PE playbooks")

    add_text_box(slide, PAD_L, Inches(1.9), Inches(5), Inches(0.3),
                 "Capital has never been this concentrated.", font_size=15, color=DARK_GREY)

    # === LEFT COLUMN (data) ===
    col_l = PAD_L
    col_w = Inches(8.0)

    # 2024 row label
    add_text_box(slide, col_l, Inches(2.35), col_w, Inches(0.25),
                 "2024 - Top 30 firms captured 75% of all US venture capital",
                 font_size=10, bold=True, color=GREY)

    # 2024 blocks
    b1_w = Inches(3.5)
    b2_w = Inches(2.2)
    b3_w = Inches(1.9)
    bh = Inches(0.75)
    by = Inches(2.65)

    s1 = add_rounded_rect(slide, col_l, by, b1_w, bh, fill_color=BLACK)
    add_text_box(slide, col_l + Inches(0.1), by + Inches(0.08), b1_w, Inches(0.35), "46%", font_size=26, bold=True, color=WHITE)
    add_text_box(slide, col_l + Inches(0.1), by + Inches(0.45), b1_w, Inches(0.25), "Top 9 firms", font_size=9, color=RGBColor(0xAA,0xAA,0xAA))

    s2 = add_rounded_rect(slide, col_l + b1_w + Inches(0.05), by, b2_w, bh, fill_color=GREY)
    add_text_box(slide, col_l + b1_w + Inches(0.15), by + Inches(0.08), b2_w, Inches(0.35), "29%", font_size=26, bold=True, color=WHITE)
    add_text_box(slide, col_l + b1_w + Inches(0.15), by + Inches(0.45), b2_w, Inches(0.25), "Next 21 firms", font_size=9, color=WHITE)

    s3 = add_rounded_rect(slide, col_l + b1_w + b2_w + Inches(0.1), by, b3_w, bh, fill_color=BG_GREY, line_color=BORDER_GREY)
    add_text_box(slide, col_l + b1_w + b2_w + Inches(0.2), by + Inches(0.08), b3_w, Inches(0.35), "25%", font_size=20, bold=True, color=GREY)
    add_text_box(slide, col_l + b1_w + b2_w + Inches(0.2), by + Inches(0.45), b3_w, Inches(0.25), "3,400+ firms", font_size=9, color=GREY)

    # 2025 row
    add_text_box(slide, col_l, Inches(3.55), col_w, Inches(0.25),
                 "2025 - It got worse", font_size=10, bold=True, color=GREY)

    r2y = Inches(3.85)
    s4 = add_rounded_rect(slide, col_l, r2y, Inches(4.5), Inches(0.7), fill_color=BLACK)
    add_text_box(slide, col_l + Inches(0.1), r2y + Inches(0.05), Inches(2), Inches(0.35), "50%+", font_size=26, bold=True, color=WHITE)
    add_text_box(slide, col_l + Inches(0.1), r2y + Inches(0.4), Inches(2), Inches(0.25), "Just 12 firms", font_size=9, color=LIGHT_GREY)

    s5 = add_rounded_rect(slide, col_l + Inches(4.55), r2y, Inches(3.0), Inches(0.7), fill_color=BG_GREY, line_color=BORDER_GREY)
    add_text_box(slide, col_l + Inches(4.65), r2y + Inches(0.05), Inches(2), Inches(0.35), "rest", font_size=20, color=LIGHT_GREY)
    add_text_box(slide, col_l + Inches(4.65), r2y + Inches(0.4), Inches(2), Inches(0.25), "Everyone else", font_size=9, color=GREY)

    # Where the money went
    add_text_box(slide, col_l, Inches(4.7), col_w, Inches(0.25),
                 "Where the money went", font_size=10, bold=True, color=GREY)

    r3y = Inches(5.0)
    ai_block = add_rounded_rect(slide, col_l, r3y, Inches(4.5), Inches(0.7), fill_color=BLACK)
    add_text_box(slide, col_l + Inches(0.1), r3y + Inches(0.05), Inches(2), Inches(0.35), "~50%", font_size=26, bold=True, color=WHITE)
    add_text_box(slide, col_l + Inches(0.1), r3y + Inches(0.4), Inches(3), Inches(0.25), "into a single category: AI", font_size=9, color=LIGHT_GREY)

    mega_block = add_rounded_rect(slide, col_l + Inches(4.55), r3y, Inches(2.5), Inches(0.7), fill_color=RGBColor(0xF5,0xF5,0xF3), line_color=BLUE)
    add_text_box(slide, col_l + Inches(4.65), r3y + Inches(0.05), Inches(2), Inches(0.35), "33%", font_size=26, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, col_l + Inches(4.65), r3y + Inches(0.4), Inches(2.3), Inches(0.25), "into 68 mega-rounds ($500M+)", font_size=8, color=GREY, alignment=PP_ALIGN.CENTER)

    # Quote
    add_text_box(slide, col_l, Inches(5.9), Inches(7), Inches(0.35),
                 '"We basically take venture risk for PE returns. I do not think any of those funds will perform as well as anyone hopes."',
                 font_size=11, color=DARK_GREY)
    add_text_box(slide, col_l, Inches(6.25), Inches(5), Inches(0.2),
                 "Head of Allocations, Vanderbilt University Endowment",
                 font_size=9, color=LIGHT_GREY)

    # === RIGHT COLUMN (cycle diagram - text representation) ===
    cx = Inches(9.5)
    cw = Inches(3.0)

    add_text_box(slide, cx, Inches(2.5), cw, Inches(0.3),
                 "THE VENTURE CONSENSUS TRAP", font_size=10, bold=True, color=LIGHT_GREY,
                 alignment=PP_ALIGN.CENTER)

    # Cycle nodes as rounded rects with text
    node_w = Inches(2.0)
    node_h = Inches(0.45)
    nc = cx + Inches(0.5)

    n1 = add_rounded_rect(slide, nc, Inches(3.0), node_w, node_h, fill_color=BLACK)
    add_text_box(slide, nc, Inches(3.02), node_w, node_h, "Capital Concentrates", font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, nc + Inches(1.8), Inches(3.5), Inches(0.5), Inches(0.3), "↓", font_size=16, color=GREY)

    n2 = add_rounded_rect(slide, nc + Inches(1.0), Inches(3.8), node_w, node_h, fill_color=GREY)
    add_text_box(slide, nc + Inches(1.0), Inches(3.82), node_w, node_h, "Consensus Grows", font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, nc + Inches(1.8), Inches(4.3), Inches(0.5), Inches(0.3), "↓", font_size=16, color=GREY)

    n3 = add_rounded_rect(slide, nc, Inches(4.6), node_w, node_h, fill_color=BLUE)
    add_text_box(slide, nc, Inches(4.62), node_w, node_h, "Prices Inflate", font_size=11, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, nc - Inches(0.3), Inches(4.3), Inches(0.5), Inches(0.3), "↑", font_size=16, color=GREY)

    n4 = add_rounded_rect(slide, nc - Inches(1.0), Inches(3.8), node_w, node_h, fill_color=GREEN)
    add_text_box(slide, nc - Inches(1.0), Inches(3.82), node_w, node_h, "Returns Compress", font_size=11, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, nc - Inches(0.3), Inches(3.5), Inches(0.5), Inches(0.3), "↑", font_size=16, color=GREY)

    # Source
    add_text_box(slide, PAD_L, Inches(6.8), Inches(8), Inches(0.2),
                 "Source: PitchBook-NVCA 2024/2025; Crunchbase 2025; MIT Strategy Science, 2024",
                 font_size=8, color=LIGHT_GREY)


def slide_3_alpha(prs):
    """Contrarian investing drives returns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    slide_title(slide, "Contrarian investing in hard domains\ndrives returns")

    # MIT callout
    mit = add_rounded_rect(slide, PAD_L, Inches(1.9), Inches(8), Inches(0.55),
                            fill_color=RGBColor(0xEE, 0xF0, 0xF5), line_color=BLUE)
    add_multi_text(slide, PAD_L + Inches(0.15), Inches(1.95), Inches(7.5), Inches(0.45), [
        ("MIT 2024: ", True, BLACK),
        ("Higher disagreement among evaluators correlates with better startup outcomes. Unique ideas spark disagreement. Common ideas spark consensus.", False, DARK_GREY),
    ], font_size=11)

    # Cycle diagram (text representation)
    cx = PAD_L + Inches(0.5)
    add_text_box(slide, cx, Inches(3.0), Inches(3), Inches(0.3),
                 "THE INFLECTION ADVANTAGE", font_size=10, bold=True, color=GREEN,
                 alignment=PP_ALIGN.CENTER)

    nc = cx + Inches(0.2)
    node_w = Inches(2.0)
    node_h = Inches(0.45)

    add_rounded_rect(slide, nc + Inches(0.3), Inches(3.5), node_w, node_h, fill_color=BLACK)
    add_text_box(slide, nc + Inches(0.3), Inches(3.52), node_w, node_h, "Hard Domains", font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, nc + Inches(2.0), Inches(4.0), Inches(0.5), Inches(0.3), "↓", font_size=16, color=GREEN)

    add_rounded_rect(slide, nc + Inches(1.2), Inches(4.3), node_w, node_h, fill_color=GREY)
    add_text_box(slide, nc + Inches(1.2), Inches(4.32), node_w, node_h, "Low Competition", font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, nc + Inches(2.0), Inches(4.8), Inches(0.5), Inches(0.3), "↓", font_size=16, color=GREEN)

    add_rounded_rect(slide, nc + Inches(0.3), Inches(5.1), node_w, node_h, fill_color=BLUE)
    add_text_box(slide, nc + Inches(0.3), Inches(5.12), node_w, node_h, "Low Entry Prices", font_size=11, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, nc - Inches(0.5), Inches(4.8), Inches(0.5), Inches(0.3), "↑", font_size=16, color=GREEN)

    add_rounded_rect(slide, nc - Inches(0.7), Inches(4.3), node_w, node_h, fill_color=GREEN)
    add_text_box(slide, nc - Inches(0.7), Inches(4.32), node_w, node_h, "High Convexity", font_size=11, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, nc - Inches(0.5), Inches(4.0), Inches(0.5), Inches(0.3), "↑", font_size=16, color=GREEN)

    # Right statement
    rx = Inches(6.5)
    add_text_box(slide, rx, Inches(3.2), Inches(5.5), Inches(1.0),
                 "Investing into hard domains is not a risk factor. It's the alpha.",
                 font_size=28, bold=True, color=BLACK)
    add_text_box(slide, rx, Inches(4.5), Inches(5.5), Inches(1.0),
                 "Hard domains are hard to access, hard to evaluate, hard to underwrite.\nTherefore competition is low and so are prices.\nThis increases convexity on a portfolio level.",
                 font_size=14, color=DARK_GREY, line_spacing=1.5)


def slide_4_opportunity(prs):
    """Diffusion table - historical + current."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    slide_title(slide, "The hard futures opportunity")
    add_text_box(slide, PAD_L, Inches(1.7), Inches(10), Inches(0.3),
                 "EUR 900B in sovereignty spending is a catalyst. The far larger opportunity is the commercial innovation it accelerates.",
                 font_size=13, color=DARK_GREY)

    # Historical column
    hist_x = PAD_L
    hist_w = Inches(5.8)
    add_text_box(slide, hist_x, Inches(2.2), hist_w, Inches(0.3),
                 "HISTORICAL: MILITARY TO CIVIL DIFFUSION", font_size=9, bold=True, color=GREY)

    historical = [
        ("Steam / Railways", "~1830", "~1930", "100x", "$2B → $200B", 100),
        ("Nuclear", "1945", "~1985", "200x", "$2B → $400B", 55),
        ("Semiconductors", "1947", "~1980", "140x", "$5B → $700B", 40),
        ("Internet", "1969", "~1995", "6000x", "$0.5B → $3T", 28),
        ("GPS / Satellites", "1978", "~1995", "10x", "$10B → $100B", 20),
        ("Drones / UAV", "2001", "~2013", "16x", "$5B → $80B", 12),
    ]

    y = Inches(2.55)
    for label, start, end, mult, expand, bar_pct in historical:
        # Label
        add_text_box(slide, hist_x, y, Inches(1.5), Inches(0.22), label, font_size=10, bold=True, color=BLACK)
        # Bar
        bar_x = hist_x + Inches(1.6)
        bar_full = Inches(2.8)
        bar_w = Inches(2.8 * bar_pct / 100)
        bar = add_rounded_rect(slide, bar_x, y + Inches(0.02), bar_w, Inches(0.18), fill_color=BORDER_GREY)
        # Year labels
        add_text_box(slide, bar_x, y - Inches(0.12), Inches(0.5), Inches(0.15), start, font_size=7, color=LIGHT_GREY)
        add_text_box(slide, bar_x + bar_w - Inches(0.4), y - Inches(0.12), Inches(0.5), Inches(0.15), end, font_size=7, color=LIGHT_GREY, alignment=PP_ALIGN.RIGHT)
        # Multiple
        add_text_box(slide, hist_x + Inches(4.6), y, Inches(0.8), Inches(0.22), mult, font_size=12, bold=True, color=GREY, alignment=PP_ALIGN.RIGHT)
        # Expansion
        add_text_box(slide, hist_x + Inches(1.6), y + Inches(0.2), Inches(3), Inches(0.15), expand, font_size=8, color=LIGHT_GREY)
        y += Inches(0.5)

    # Current wave column
    curr_x = Inches(7.0)
    curr_w = Inches(5.8)
    add_text_box(slide, curr_x, Inches(2.2), curr_w, Inches(0.3),
                 "NOW: DIFFUSION LAG APPROACHING ZERO", font_size=9, bold=True, color=GREY)

    current = [
        ("Space Economy", "2010", "2040", "55x", "$33B → $1.8T", 100),
        ("Edge AI", "2015", "2035", "35x", "$10B → $350B", 82),
        ("Autonomous Robotics", "2010", "2035", "50x", "$20B → $1T", 73),
        ("Hypersonics", "2015", "2035", "43x", "$7B → $300B", 60),
        ("Distributed Mfg", "2018", "2035", "60x", "$5B → $300B", 50),
        ("Frontier Compute", "2020", "2040", "250x", "$2B → $500B", 40),
    ]

    y = Inches(2.55)
    for label, start, end, mult, expand, bar_pct in current:
        add_text_box(slide, curr_x, y, Inches(1.8), Inches(0.22), label, font_size=10, bold=True, color=BLACK)
        bar_x = curr_x + Inches(1.9)
        bar_w = Inches(2.8 * bar_pct / 100)
        bar = add_rounded_rect(slide, bar_x, y + Inches(0.02), bar_w, Inches(0.18), fill_color=GREEN)
        add_text_box(slide, bar_x, y - Inches(0.12), Inches(0.5), Inches(0.15), start, font_size=7, color=LIGHT_GREY)
        add_text_box(slide, bar_x + bar_w - Inches(0.4), y - Inches(0.12), Inches(0.5), Inches(0.15), end, font_size=7, color=LIGHT_GREY, alignment=PP_ALIGN.RIGHT)
        add_text_box(slide, curr_x + Inches(4.9), y, Inches(0.8), Inches(0.22), mult, font_size=12, bold=True, color=BLACK, alignment=PP_ALIGN.RIGHT)
        add_text_box(slide, curr_x + Inches(1.9), y + Inches(0.2), Inches(3), Inches(0.15), expand, font_size=8, color=GREEN)
        y += Inches(0.5)

    # Bottom insight
    add_text_box(slide, PAD_L, Inches(5.8), Inches(11), Inches(0.5),
                 "The diffusion lag from military invention to commercial market has collapsed from centuries to decades to years. Fund III invests at the moment this lag approaches zero.",
                 font_size=12, color=DARK_GREY)

    add_text_box(slide, PAD_L, Inches(6.8), Inches(10), Inches(0.2),
                 "Market sizing: McKinsey, Morgan Stanley, Teal Group, Grand View Research",
                 font_size=8, color=LIGHT_GREY)


def slide_5_portfolio(prs):
    """Portfolio grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    slide_title(slide, "We invest in hard futures at pre-seed\nwhere few can follow")

    companies = [
        ("Radical", "Perpetual flight - stratospheric platforms at 20km altitude"),
        ("Deep Earth", "Geospatial AI mapping the earth's subsurface"),
        ("Foundational", "Space laser ranging for ultra-precise geodetic data"),
        ("Fabric", "Custom silicon for encrypted computation and edge AI"),
        ("Ubitium", "Universal processor for any AI workload at the edge"),
        ("Hedy", "Network invisibility in contested environments"),
        ("NAD", "Modular air defence - Europe's Iron Dome"),
        ("Ark", "Autonomous fleet control for drones and robots"),
        ("Levtek", "General-purpose industrial robots for blue-collar workers"),
        ("Lodestar", "Bodyguard satellites - robotic protection of orbital assets"),
    ]

    cols = 5
    card_w = Inches(2.2)
    card_h = Inches(1.5)
    gap = Inches(0.25)
    start_x = PAD_L + Inches(0.2)
    start_y = Inches(2.2)

    for i, (name, desc) in enumerate(companies):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + Inches(0.2))

        # Image placeholder
        add_rounded_rect(slide, x, y, card_w, Inches(0.7), fill_color=BG_GREY, line_color=BORDER_GREY)
        add_text_box(slide, x, y + Inches(0.15), card_w, Inches(0.3), "[ image ]",
                     font_size=9, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x, y + Inches(0.75), card_w, Inches(0.25), name,
                     font_size=13, bold=True, color=BLACK)
        add_text_box(slide, x, y + Inches(1.0), card_w, Inches(0.45), desc,
                     font_size=9, color=GREY, line_spacing=1.3)

    # Properties bar
    props = [
        ("Physically hard to reach", "Stratosphere. Orbit. Underground. Contested zones.", GREEN),
        ("Technically hard to build", "Years of deep engineering before first unit ships.", BLUE),
        ("Psychologically hard to fund", "No playbooks. No Gartner TAM. Dies in committee.", BLACK),
    ]
    py = Inches(5.8)
    pw = Inches(3.6)
    for i, (label, sub, color) in enumerate(props):
        px = PAD_L + Inches(0.2) + i * (pw + Inches(0.3))
        add_rounded_rect(slide, px, py, pw, Inches(0.65), fill_color=None, line_color=color)
        add_text_box(slide, px + Inches(0.1), py + Inches(0.05), pw - Inches(0.2), Inches(0.25),
                     label, font_size=11, bold=True, color=color)
        add_text_box(slide, px + Inches(0.1), py + Inches(0.3), pw - Inches(0.2), Inches(0.3),
                     sub, font_size=9, color=GREY)


def slide_6_timeline(prs):
    """Macro timeline - 2 cycles."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    slide_title(slide, "We systematically invest before the market")

    # Cycle 1
    y1 = Inches(1.8)
    add_rounded_rect(slide, PAD_L, y1, Inches(1.5), Inches(2.0), fill_color=BG_GREY, line_color=BORDER_GREY)
    add_text_box(slide, PAD_L + Inches(0.1), y1 + Inches(0.1), Inches(1.3), Inches(0.3), "2019", font_size=20, bold=True, color=BLACK)
    add_text_box(slide, PAD_L + Inches(0.1), y1 + Inches(0.5), Inches(1.3), Inches(0.5),
                 "Depressed crypto.\nNobody looking.", font_size=11, color=GREY)

    positions_1 = ["BTC at $3k", "ETH at $200", "Centrifuge $30M", "Anytype $5M", "Molecule $5M"]
    for i, pos in enumerate(positions_1):
        add_text_box(slide, Inches(2.8), y1 + Inches(0.15) + i * Inches(0.3), Inches(2.5), Inches(0.25),
                     pos, font_size=11, color=BLACK)

    add_text_box(slide, Inches(5.5), y1 + Inches(0.5), Inches(0.5), Inches(0.3), "→", font_size=24, color=GREEN)

    add_text_box(slide, Inches(6.2), y1 + Inches(0.2), Inches(2.5), Inches(0.3), "Crypto rally 2021", font_size=11, bold=True, color=DARK_GREEN)
    add_text_box(slide, Inches(6.2), y1 + Inches(0.5), Inches(2.5), Inches(0.3), "Positions 10-50x", font_size=11, bold=True, color=DARK_GREEN)

    add_text_box(slide, Inches(8.8), y1 + Inches(0.5), Inches(0.5), Inches(0.3), "→", font_size=24, color=GREEN)

    add_text_box(slide, Inches(9.5), y1 + Inches(0.15), Inches(1.5), Inches(0.5), "8x", font_size=36, bold=True, color=DARK_GREEN)
    add_text_box(slide, Inches(9.5), y1 + Inches(0.65), Inches(1.5), Inches(0.25), "TVPI", font_size=10, color=GREY)
    add_text_box(slide, Inches(10.8), y1 + Inches(0.15), Inches(1.5), Inches(0.5), "3x", font_size=28, bold=True, color=DARK_GREEN)
    add_text_box(slide, Inches(10.8), y1 + Inches(0.55), Inches(1.5), Inches(0.25), "DPI", font_size=10, color=GREY)

    # Cycle 2
    y2 = Inches(4.2)
    add_rounded_rect(slide, PAD_L, y2, Inches(1.5), Inches(2.0), fill_color=BG_GREY, line_color=BORDER_GREY)
    add_text_box(slide, PAD_L + Inches(0.1), y2 + Inches(0.1), Inches(1.3), Inches(0.3), "2022", font_size=20, bold=True, color=BLACK)
    add_text_box(slide, PAD_L + Inches(0.1), y2 + Inches(0.5), Inches(1.3), Inches(0.5),
                 "Ukraine invasion.\nEU sovereignty\nwake-up.", font_size=11, color=GREY)

    positions_2 = ["Radical 2023", "Ubitium 2023", "Ark 2024", "Hedy 2024", "NAD 2025"]
    for i, pos in enumerate(positions_2):
        add_text_box(slide, Inches(2.8), y2 + Inches(0.15) + i * Inches(0.3), Inches(2.5), Inches(0.25),
                     pos, font_size=11, color=BLACK)

    add_text_box(slide, Inches(5.5), y2 + Inches(0.5), Inches(0.5), Inches(0.3), "→", font_size=24, color=GREEN)

    validations = ["EUR 800B+ committed", "SAFE disbursements 2026", "65% EU content rule"]
    for i, v in enumerate(validations):
        add_text_box(slide, Inches(6.2), y2 + Inches(0.2) + i * Inches(0.3), Inches(2.5), Inches(0.3),
                     v, font_size=11, bold=True, color=DARK_GREEN)

    add_text_box(slide, Inches(8.8), y2 + Inches(0.5), Inches(0.5), Inches(0.3), "→", font_size=24, color=GREEN)
    add_text_box(slide, Inches(9.5), y2 + Inches(0.15), Inches(2), Inches(0.5), "TBD", font_size=36, bold=True, color=GREEN)
    add_text_box(slide, Inches(9.5), y2 + Inches(0.65), Inches(2.5), Inches(0.25), "Procurement wave arriving", font_size=10, color=GREY)

    # Bottom
    add_text_box(slide, PAD_L, Inches(6.6), Inches(10), Inches(0.4),
                 "2 cycles, same pattern. We invest when it is hard. The market catches up. We can do it again.",
                 font_size=14, bold=True, color=BLACK)


def slide_7_prices(prs):
    """Entry prices table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    slide_title(slide, "We enter at prices that create convexity")

    from pptx.util import Inches, Pt
    from pptx.table import Table as PptxTable

    rows_data = [
        ["Hedy", "Network security", "2025", "1.25 / 9.8", "12.8%", "$3M ARR; SF + enterprise; $30M pipeline", "$40M post", "4.1x", "Closing"],
        ["Ark", "Drone + robot fleet", "2024", "1.45 / 21.5", "6.8%", "$5M rev; $100M pipeline; 20+ UA brigades", "$65M post", "3.0x", "Closed"],
        ["Hanseatic", "Maritime intel", "-", "2.00 / 25", "8.0%", "$150M turnover 2026; $300M pipeline", "$60M post", "2.4x", "Term sheet"],
        ["NAD", "Modular air defence", "2025", "1.50 / 25", "6.0%", "Working prototype; $100M pipeline", "$56M+", "2.2x", "Equity imm."],
        ["Levtek", "Industrial robots", "2025", "1.18 / 12.25", "8.2%", "5 units IKEA/Sony; $500K rev Y1", "-", "-", "Deployed"],
        ["Radical", "Stratospheric flight", "2023", "1.00 / 17.5", "4.5%", "Virgin flight Q3 2025; $50M+ LOIs", "-", "-", "Raising"],
        ["Foundational", "Space laser ranging", "2025", "1.50 / 17", "8.0%", "$5M revenue; pilot station built", "$25M post", "1.5x", "Closed"],
        ["Ubitium", "Edge processor", "2024", "1.10 / 14", "8.0%", "2 chip tape-outs in 2026", "-", "-", "Deployed"],
        ["Lodestar", "Satellite servicing", "2024", "1.77 / 16", "9.8%", "$5M rev 2026; UK→US; space launch", "-", "-", "Raising"],
        ["Deep Earth", "Subsurface AI", "2024", "1.00 / 9.5", "7.2%", "$1M ARR in Y2", "$20M (est)", "2.1x", "Deployed"],
        ["Fabric", "Encrypted silicon", "2023", "2.00 / 55", "2.5%", "$50M+ LOIs; chip taped out", "-", "-", "Deployed"],
    ]

    headers = ["Company", "Domain", "Year", "Inv / Entry", "Own%", "Fundamentals", "Current", "Mult", "Status"]
    col_widths = [Inches(1.1), Inches(1.3), Inches(0.5), Inches(1.0), Inches(0.6), Inches(3.8), Inches(1.0), Inches(0.6), Inches(0.8)]

    n_rows = len(rows_data) + 1
    n_cols = len(headers)

    table_shape = slide.shapes.add_table(n_rows, n_cols, PAD_L, Inches(1.7), sum(col_widths, Emu(0)), Inches(4.8))
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Style header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = GREY
            p.font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE

    # Fill data
    for i, row in enumerate(rows_data):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9) if j != 5 else Pt(8)  # Fundamentals smaller
                p.font.name = FONT
                p.font.color.rgb = BLACK
                if j == 0:  # Company name bold
                    p.font.bold = True
                if j == 1:  # Domain grey
                    p.font.color.rgb = GREY
                if j == 7 and val not in ["-"]:  # Multiple green
                    p.font.bold = True
                    p.font.color.rgb = DARK_GREEN
                if j == 5:  # Fundamentals
                    p.font.color.rgb = DARK_GREY
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

    # Bottom note
    add_multi_text(slide, PAD_L, Inches(6.6), Inches(10), Inches(0.3), [
        ("Five markups between 1.5x and 4.1x. None reflected in the reported 0.95x TVPI.", True, BLACK),
    ], font_size=12)
    add_text_box(slide, PAD_L, Inches(6.9), Inches(8), Inches(0.3),
                 "Not because we negotiated harder. Because we were there before anyone else was looking.",
                 font_size=11, color=DARK_GREY)


def slide_8_system(prs):
    """Portfolio system - Sense → Compute → Act."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    slide_title(slide, "We build portfolios as systems")

    columns = [
        ("SENSE", "Planetary\nIntelligence", BLUE,
         [("Radical", "Stratosphere"), ("Deep Earth", "Subsurface"), ("Foundational", "Space")]),
        ("COMPUTE", "Sovereign\nCompute", GREEN,
         [("Fabric", "Encrypted silicon"), ("Ubitium", "Edge AI"), ("Hedy", "Network stealth")]),
        ("ACT", "Autonomous\nMachines", BLACK,
         [("NAD", "Air defence"), ("Ark", "Fleet control"), ("Levtek", "Industrial robots"), ("Lodestar", "Orbital servicing")]),
    ]

    for i, (verb, name, color, companies) in enumerate(columns):
        x = PAD_L + Inches(0.3) + i * Inches(4.0)
        w = Inches(3.2)

        # Header box
        add_rounded_rect(slide, x, Inches(2.0), w, Inches(1.2), fill_color=None, line_color=color)
        add_text_box(slide, x + Inches(0.15), Inches(2.05), w, Inches(0.3), verb,
                     font_size=10, bold=True, color=color)
        add_text_box(slide, x + Inches(0.15), Inches(2.35), w, Inches(0.6), name,
                     font_size=20, bold=True, color=BLACK)

        # Arrow between columns
        if i < 2:
            ax = x + w + Inches(0.15)
            add_text_box(slide, ax, Inches(2.4), Inches(0.5), Inches(0.4), "→",
                         font_size=24, bold=True, color=color)

        # Companies
        for j, (co_name, co_role) in enumerate(companies):
            cy = Inches(3.4) + j * Inches(0.4)
            add_multi_text(slide, x + Inches(0.15), cy, w, Inches(0.35), [
                (co_name, True, BLACK),
                (f"  {co_role}", False, GREY),
            ], font_size=12)

    # Bottom
    add_text_box(slide, PAD_L, Inches(5.5), Inches(11), Inches(0.3),
                 "Intelligence feeds compute  →  Compute drives machines  →  The portfolio is a closed loop",
                 font_size=14, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, PAD_L, Inches(5.9), Inches(11), Inches(0.3),
                 "The portfolio is highly synergetic. We never back competing companies.",
                 font_size=12, color=GREY, alignment=PP_ALIGN.CENTER)


def slide_9_team(prs):
    """Team slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    slide_title(slide, "We built through 7y of compounding conviction")

    add_text_box(slide, PAD_L, Inches(1.7), Inches(10), Inches(0.4),
                 "We are engineers, operators and investors by training. We have been through hardship and are as resilient as our founders. We collaborate since 2019 (some of us since the 1990s).",
                 font_size=13, color=DARK_GREY)

    team = [
        ("Alexander Lange", "Founding GP", "Investment Team",
         "Law + Econ; top 5%; KAS fellow (top 1%). Google, fintech. 5y at Index Ventures + Earlybird. Raised $40M as Solo GP through a pandemic."),
        ("Jonatan Luther-Bergquist", "GP", "",
         "Engineering Physicist, UU/EPFL/TUM. Top 1% in SWE aptitude tests. Published researcher. 3y UN product (500k users). 2y BCG, IT architecture."),
        ("Alexander Patow", "Engineering", "Engineering",
         "Mechanical engineer, BU. 4y Accenture SWE. 4y EQT / Motherbrain; founding data analytics. Board Stockholm AI. Built Inflection's entire tech stack."),
    ]

    for i, (name, role, label, bio) in enumerate(team):
        x = PAD_L + i * Inches(4.0)
        y = Inches(2.3)
        w = Inches(3.6)

        if label:
            add_text_box(slide, x, y, w, Inches(0.2), label.upper(),
                         font_size=8, bold=True, color=GREEN if i < 2 else BLUE)

        # Photo placeholder
        add_rounded_rect(slide, x, y + Inches(0.25), w, Inches(0.8), fill_color=BG_GREY, line_color=BORDER_GREY)
        add_text_box(slide, x, y + Inches(0.45), w, Inches(0.3), "[ photo ]",
                     font_size=9, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x, y + Inches(1.15), w, Inches(0.25), name, font_size=13, bold=True, color=BLACK)
        add_text_box(slide, x, y + Inches(1.35), w, Inches(0.2), role, font_size=10, bold=True, color=GREY)
        add_text_box(slide, x, y + Inches(1.55), w, Inches(0.8), bio, font_size=9, color=GREY, line_spacing=1.4)

    # Ops row
    ops = [
        ("Rebecca Mahoney", "VP Finance"),
        ("Robert Shapiro", "Asset & Risk Mgmt"),
        ("Jonathan Levin", "General Counsel"),
    ]
    oy = Inches(5.0)
    add_text_box(slide, PAD_L, oy, Inches(2), Inches(0.2), "OPERATIONS",
                 font_size=8, bold=True, color=GREY)
    for i, (name, role) in enumerate(ops):
        x = PAD_L + i * Inches(4.0)
        add_rounded_rect(slide, x, oy + Inches(0.25), Inches(3.6), Inches(0.6), fill_color=BG_GREY, line_color=BORDER_GREY)
        add_text_box(slide, x, oy + Inches(0.35), Inches(3.6), Inches(0.2), "[ photo ]",
                     font_size=9, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, oy + Inches(0.9), Inches(3.6), Inches(0.2), name, font_size=11, bold=True, color=BLACK)
        add_text_box(slide, x, oy + Inches(1.1), Inches(3.6), Inches(0.2), role, font_size=9, color=GREY)

    # Logo bar
    logos = ["Index Ventures", "Earlybird", "Google", "TUM", "BCG", "EPFL", "EQT"]
    ly = Inches(6.6)
    for i, logo in enumerate(logos):
        lx = PAD_L + i * Inches(1.7)
        add_text_box(slide, lx, ly, Inches(1.5), Inches(0.3), logo,
                     font_size=12, bold=True, color=GREY)


def slide_10_machine(prs):
    """Two machines."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    slide_title(slide, "Two machines to punch above our weight")

    cards = [
        ("Dealflow + Portfolio Platform", "Matter... / EDTH",
         "We co-founded a platform that runs hackathons in defence, energy, future of computing, aerospace, bio - with a media arm. Beyond dealflow: engineering hires, BD, procurement access, portfolio support.",
         [("25", "Hackathons / year"), ("100s", "Founders scouted"), ("5+", "Domains covered")]),
        ("Operations Engine", "CAVI",
         "Computer Aided Venture Investing. Automated deal flow scoring, portfolio monitoring, and LP reporting. Three people running what most funds need nine for.",
         [("3", "Team members"), ("9", "Equivalent headcount"), ("3x", "Leverage ratio")]),
    ]

    for i, (badge, title, body, stats) in enumerate(cards):
        x = PAD_L + i * Inches(6.2)
        w = Inches(5.6)
        y = Inches(1.8)

        add_rounded_rect(slide, x, y, w, Inches(4.5), fill_color=None, line_color=BORDER_GREY)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3), Inches(0.2),
                     badge.upper(), font_size=8, bold=True, color=GREY)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.4), w - Inches(0.4), Inches(0.4),
                     title, font_size=22, bold=True, color=BLACK)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.9), w - Inches(0.4), Inches(1.2),
                     body, font_size=12, color=DARK_GREY, line_spacing=1.5)

        # Image placeholder
        add_rounded_rect(slide, x + Inches(0.2), y + Inches(2.2), w - Inches(0.4), Inches(1.0),
                         fill_color=BG_GREY, line_color=BORDER_GREY)
        add_text_box(slide, x + Inches(0.2), y + Inches(2.5), w - Inches(0.4), Inches(0.3),
                     "[ image ]", font_size=9, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

        # Stats
        for j, (num, label) in enumerate(stats):
            sx = x + Inches(0.2) + j * Inches(1.7)
            add_text_box(slide, sx, y + Inches(3.4), Inches(1.5), Inches(0.4),
                         num, font_size=24, bold=True, color=BLACK)
            add_text_box(slide, sx, y + Inches(3.8), Inches(1.5), Inches(0.25),
                         label, font_size=9, color=GREY)

    add_text_box(slide, PAD_L, Inches(6.5), Inches(10), Inches(0.3),
                 "Direct pipeline into exactly our domains. No cold outreach. No spray and pray.",
                 font_size=13, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)


def slide_11_numbers(prs):
    """Honest numbers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    slide_title(slide, "Our numbers - the good, the bad, the ugly")

    # Fund I card
    f1x = PAD_L
    f1w = Inches(5.5)
    add_rounded_rect(slide, f1x, Inches(1.8), f1w, Inches(1.8), fill_color=None, line_color=BORDER_GREY)
    add_text_box(slide, f1x + Inches(0.2), Inches(1.9), Inches(2), Inches(0.3), "Fund I", font_size=18, bold=True, color=BLACK)
    add_text_box(slide, f1x + Inches(0.2), Inches(2.2), Inches(3), Inches(0.2), "$1.1M | 2019 vintage", font_size=11, color=GREY)
    add_text_box(slide, f1x + Inches(0.2), Inches(2.5), Inches(1.5), Inches(0.5), "8.0x", font_size=40, bold=True, color=DARK_GREEN)
    add_text_box(slide, f1x + Inches(2.0), Inches(2.6), Inches(1.5), Inches(0.4), "3.0x DPI", font_size=24, bold=True, color=BLACK)
    add_text_box(slide, f1x + Inches(0.2), Inches(3.0), f1w - Inches(0.4), Inches(0.4),
                 "Proof of concept. Pandemic cut fundraise from $10M to $1.1M. Selection engine validated. 5 of 12 positions drove returns.",
                 font_size=10, color=GREY)

    # Fund II card
    f2x = Inches(7.0)
    f2w = Inches(5.5)
    add_rounded_rect(slide, f2x, Inches(1.8), f2w, Inches(1.8), fill_color=None, line_color=BORDER_GREY)
    add_text_box(slide, f2x + Inches(0.2), Inches(1.9), Inches(2), Inches(0.3), "Fund II", font_size=18, bold=True, color=BLACK)
    add_text_box(slide, f2x + Inches(0.2), Inches(2.2), Inches(3), Inches(0.2), "$40M | 2022 vintage", font_size=11, color=GREY)
    add_text_box(slide, f2x + Inches(0.2), Inches(2.5), Inches(2), Inches(0.5), "0.95x", font_size=40, bold=True, color=BLACK)
    add_text_box(slide, f2x + Inches(0.2), Inches(2.95), Inches(3), Inches(0.2), "TVPI (most conservative basis)", font_size=10, color=GREY)
    add_text_box(slide, f2x + Inches(0.2), Inches(3.15), f2w - Inches(0.4), Inches(0.3),
                 "Below water. Most aggressive write-downs, most conservative markups. No SAFE uplifts.",
                 font_size=10, color=GREY)

    # Slope section
    add_text_box(slide, PAD_L, Inches(3.9), Inches(8), Inches(0.25),
                 "THE SLOPE BEHIND 0.95X - MARKUPS NOT YET IN REPORTED TVPI",
                 font_size=9, bold=True, color=GREY)

    markups = [
        ("Hedy", "$9.8M → $40M", "4.1x"),
        ("Ark", "$21.5M → $65M", "3.0x"),
        ("Hanseatic", "$25M → $60M", "2.4x"),
        ("NAD", "$25M → $56M+", "2.2x"),
        ("Foundational", "$17M → $25M", "1.5x"),
    ]

    sx = PAD_L
    sw = Inches(2.2)
    for i, (name, detail, mult) in enumerate(markups):
        mx = sx + i * (sw + Inches(0.2))
        my = Inches(4.3)
        add_rounded_rect(slide, mx, my, sw, Inches(0.9), fill_color=BG_GREY, line_color=BORDER_GREY)
        add_text_box(slide, mx + Inches(0.1), my + Inches(0.05), sw, Inches(0.25), name, font_size=12, bold=True, color=BLACK)
        add_text_box(slide, mx + Inches(0.1), my + Inches(0.3), sw, Inches(0.2), detail, font_size=10, color=DARK_GREY)
        add_text_box(slide, mx + Inches(0.1), my + Inches(0.55), sw, Inches(0.25), mult, font_size=18, bold=True, color=DARK_GREEN)

    add_multi_text(slide, PAD_L, Inches(5.5), Inches(11), Inches(0.3), [
        ("a16z founders (Isomer, Hutt Capital, JT Esquibel) saw these numbers and recommitted to Fund III. ", True, BLACK),
        ("Partners at General Catalyst, Point Nine, Earlybird, and Crane also committed personally to Fund II.", False, DARK_GREY),
    ], font_size=11)


def slide_12_mars(prs):
    """Fund III - Mars."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    slide_title(slide, "Fund III - Mars")

    metrics = [
        ("$50M", "Target fund size", True),
        ("24", "Target positions", False),
        ("6-10x", "Target return multiple", True),
        ("10y", "Fund life cycle", False),
    ]

    for i, (num, label, highlight) in enumerate(metrics):
        mx = PAD_L + i * Inches(3.0)
        my = Inches(1.9)
        mw = Inches(2.6)
        bg = RGBColor(0xF5, 0xF5, 0xF3) if highlight else None
        border = BORDER_GREY if not highlight else None
        add_rounded_rect(slide, mx, my, mw, Inches(1.0), fill_color=bg, line_color=border)
        add_text_box(slide, mx + Inches(0.15), my + Inches(0.1), mw, Inches(0.5), num,
                     font_size=32, bold=True, color=BLACK)
        add_text_box(slide, mx + Inches(0.15), my + Inches(0.6), mw, Inches(0.3), label,
                     font_size=11, color=GREY)

    # Timeline bar
    ty = Inches(3.3)
    add_rounded_rect(slide, PAD_L, ty, Inches(4.8), Inches(0.35), fill_color=BLACK)
    add_text_box(slide, PAD_L + Inches(0.1), ty + Inches(0.02), Inches(4.6), Inches(0.3),
                 "4y invest period", font_size=10, bold=True, color=WHITE)

    add_rounded_rect(slide, PAD_L + Inches(4.85), ty, Inches(4.8), Inches(0.35), fill_color=GREY)
    add_text_box(slide, PAD_L + Inches(4.95), ty + Inches(0.02), Inches(4.6), Inches(0.3),
                 "Harvest", font_size=10, bold=True, color=WHITE)

    add_rounded_rect(slide, PAD_L + Inches(9.7), ty, Inches(2.4), Inches(0.35), fill_color=BG_GREY, line_color=BORDER_GREY)
    add_text_box(slide, PAD_L + Inches(9.8), ty + Inches(0.02), Inches(2.2), Inches(0.3),
                 "2+2y extension", font_size=10, color=GREY)

    # Capital allocation
    add_text_box(slide, PAD_L, Inches(4.0), Inches(3), Inches(0.2),
                 "CAPITAL ALLOCATION", font_size=9, bold=True, color=GREY)
    add_rounded_rect(slide, PAD_L, Inches(4.3), Inches(9.5), Inches(0.35), fill_color=BLACK)
    add_text_box(slide, PAD_L + Inches(0.1), Inches(4.32), Inches(5), Inches(0.3),
                 "80% Initial investments", font_size=10, bold=True, color=WHITE)
    add_rounded_rect(slide, PAD_L + Inches(9.55), Inches(4.3), Inches(2.55), Inches(0.35), fill_color=GREEN)
    add_text_box(slide, PAD_L + Inches(9.65), Inches(4.32), Inches(2.4), Inches(0.3),
                 "20% Follow-on", font_size=10, bold=True, color=BLACK)

    # Closing timeline
    items = [
        ("First close - Mid 2025", "$6M committed", "Isomer, Hutt Capital, JT Esquibel (a16z founders), Christian Reber"),
        ("Final close", "December 2026", ""),
        ("Terms", "20% carry | 2% mgmt fee p.a.", ""),
    ]
    iy = Inches(5.2)
    for i, (label, value, sub) in enumerate(items):
        ix = PAD_L + i * Inches(4.0)
        iw = Inches(3.6)
        add_rounded_rect(slide, ix, iy, iw, Inches(1.2) if sub else Inches(0.8),
                         fill_color=BG_GREY if i == 0 else None, line_color=BORDER_GREY)
        add_text_box(slide, ix + Inches(0.15), iy + Inches(0.05), iw, Inches(0.2),
                     label, font_size=9, bold=True, color=GREY)
        add_text_box(slide, ix + Inches(0.15), iy + Inches(0.3), iw, Inches(0.3),
                     value, font_size=14, bold=True, color=BLACK)
        if sub:
            add_text_box(slide, ix + Inches(0.15), iy + Inches(0.65), iw - Inches(0.3), Inches(0.4),
                         sub, font_size=9, color=GREY)


def slide_13_beliefs(prs):
    """What you need to believe."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)

    add_text_box(slide, Inches(0), Inches(1.5), SLIDE_W, Inches(0.8),
                 "What you need to believe", font_size=36, bold=True, color=BLACK,
                 alignment=PP_ALIGN.CENTER)

    beliefs = [
        ("1", "Hard domains produce better returns than consensus domains. Difficulty filters competition and compresses entry prices."),
        ("2", "European frontier tech is a multi-decade structural shift, not a political cycle."),
        ("3", "A small, thematic team with 7 years of consistency outperforms a large platform deploying across twenty sectors."),
    ]

    for i, (num, text) in enumerate(beliefs):
        by = Inches(2.8) + i * Inches(1.2)
        bx = Inches(3.0)
        bw = Inches(7.5)

        add_rounded_rect(slide, bx, by, bw, Inches(0.9), fill_color=None, line_color=BORDER_GREY)
        add_text_box(slide, bx + Inches(0.2), by + Inches(0.1), Inches(0.5), Inches(0.5),
                     num, font_size=28, bold=True, color=GREEN)
        add_text_box(slide, bx + Inches(0.7), by + Inches(0.15), bw - Inches(1.0), Inches(0.6),
                     text, font_size=14, color=BLACK, line_spacing=1.4)

    add_text_box(slide, Inches(0), Inches(6.3), SLIDE_W, Inches(0.3),
                 "Our fundraising period ends on December 31, 2026.",
                 font_size=13, color=GREY, alignment=PP_ALIGN.CENTER)


def slide_14_references(prs):
    """References."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    slide_title(slide, "References")

    # Left column - Quotes
    lx = PAD_L
    lw = Inches(5.5)

    add_text_box(slide, lx, Inches(1.7), lw, Inches(0.2),
                 "WHAT FOUNDERS AND LPS SAY", font_size=8, bold=True, color=GREY)

    quotes = [
        ('"You guys went repeatedly into territory that is immensely painful. That will be rewarded."',
         "Manuel Grossmann, Armino Collective"),
        ('"[Placeholder: founder quote about conviction or value-add]"', "[Founder Name], [Company]"),
        ('"[Placeholder: LP quote about differentiation]"', "[LP Name], [Organization]"),
    ]
    qy = Inches(2.0)
    for text, attr in quotes:
        add_text_box(slide, lx, qy, lw, Inches(0.35), text, font_size=11, color=BLACK)
        add_text_box(slide, lx, qy + Inches(0.35), lw, Inches(0.2), attr, font_size=9, color=GREY)
        qy += Inches(0.7)

    add_text_box(slide, lx, qy + Inches(0.1), lw, Inches(0.2),
                 "LPS AND BACKERS", font_size=8, bold=True, color=GREY)
    brands = "Isomer Capital  |  Hutt Capital  |  JT Esquibel  |  Christian Reber  |  General Catalyst  |  Point Nine  |  Earlybird  |  Crane"
    add_text_box(slide, lx, qy + Inches(0.35), lw, Inches(0.4), brands, font_size=10, color=BLACK)

    # Right column - Writing
    rx = Inches(7.0)
    rw = Inches(5.5)

    add_text_box(slide, rx, Inches(1.7), rw, Inches(0.2),
                 "SELECTED WRITING AND MEDIA", font_size=8, bold=True, color=GREY)

    articles = [
        ("Heresy and the Venture Industrial Complex", "Why non-consensus investing produces better returns"),
        ("Europe's New Defense - The Investment Thesis", "EUR 900B in sovereignty spending and deep tech"),
        ("Depth vs Breadth in Venture Capital", "Why thematic focus beats generalist deployment"),
        ("The Anatomy of Inflections", "Identifying phase transitions before consensus"),
        ("Matter.../EDTH - Building the Ecosystem", "Hackathons and media across defense, energy, aerospace"),
        ("The Sovereign Compute Stack", "Why Europe needs its own hardware-software supply chain"),
        ("Why We Back First-Time Founders", "Conviction-driven investing at the earliest stage"),
    ]
    ay = Inches(2.0)
    for title, desc in articles:
        add_text_box(slide, rx, ay, rw, Inches(0.2), title, font_size=11, bold=True, color=BLACK)
        add_text_box(slide, rx, ay + Inches(0.2), rw, Inches(0.2), desc, font_size=9, color=GREY)
        ay += Inches(0.45)

    add_text_box(slide, rx, ay + Inches(0.1), rw, Inches(0.3),
                 "inflection.xyz/blog", font_size=12, bold=True, color=BLUE)


def slide_15_endcard(prs):
    """Closing card."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BLACK

    add_text_box(slide, Inches(0), Inches(2.0), SLIDE_W, Inches(1.0),
                 "inflection", font_size=56, bold=True, color=RGBColor(0x80, 0x80, 0x80),
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(4.0), Inches(9.3), Inches(1.0),
                 "Conviction, not conformity. Curiosity, not playbooks.\nWe reclaim venture's origins to power a European renaissance.",
                 font_size=18, color=RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.CENTER,
                 line_spacing=1.6)


def main():
    print("Building editable PPTX...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_0_title(prs)
    print("  Slide 0: Title")
    slide_1_origins(prs)
    print("  Slide 1: Origins")
    slide_2_concentration(prs)
    print("  Slide 2: Concentration")
    slide_3_alpha(prs)
    print("  Slide 3: Alpha")
    slide_4_opportunity(prs)
    print("  Slide 4: Opportunity")
    slide_5_portfolio(prs)
    print("  Slide 5: Portfolio")
    slide_6_timeline(prs)
    print("  Slide 6: Timeline")
    slide_7_prices(prs)
    print("  Slide 7: Prices table")
    slide_8_system(prs)
    print("  Slide 8: System")
    slide_9_team(prs)
    print("  Slide 9: Team")
    slide_10_machine(prs)
    print("  Slide 10: Machine")
    slide_11_numbers(prs)
    print("  Slide 11: Numbers")
    slide_12_mars(prs)
    print("  Slide 12: Mars fund")
    slide_13_beliefs(prs)
    print("  Slide 13: Beliefs")
    slide_14_references(prs)
    print("  Slide 14: References")
    slide_15_endcard(prs)
    print("  Slide 15: Endcard")

    prs.save(str(PPTX_PATH))
    print(f"\nSaved: {PPTX_PATH}")


if __name__ == "__main__":
    main()
